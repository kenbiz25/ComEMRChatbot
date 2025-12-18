
# main.py — ComEMR Support Backend (Enterprise RAG + Learning, Role-Agnostic)
# ==============================================================================
# ✔ KB-first RAG (FAQ-aware), role-agnostic (no role filtering or tone)
# ✔ Streaming via OpenAI Responses API + safe non-stream fallback
# ✔ Robust sessions & conversation JSONL persistence
# ✔ Intent logging (embeddings) + answer quality logs
# ✔ Admin learning approval & KB reindex
# ✔ Confidence scoring & source citations + confidence-based phrasing
# ✔ Dynamic threshold + fallback to top hits
# ✔ JSON + Form input support for endpoints
# ✔ Request ID logging + structured error responses
# ✔ Atomic file writes to avoid corruption (Windows-safe)
# ✔ Safe session ID generation
# ✔ /config, /health, /kb/stats endpoints
# ✔ Common fallbacks for frequent queries

# -------------------- Standard Library --------------------
import os
import json
import pickle
import uuid
import logging
import zipfile
import re
import tempfile
from datetime import datetime
from pathlib import Path
from typing import Dict, Any, Iterator, List, Optional

# -------------------- FastAPI --------------------
from fastapi import FastAPI, Request, Form, HTTPException, Body, APIRouter
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse, HTMLResponse, JSONResponse

# -------------------- Third Party --------------------
from dotenv import load_dotenv
from PyPDF2 import PdfReader
import faiss
import numpy as np
from sentence_transformers import SentenceTransformer
from openai import OpenAI
from docx import Document
from pptx import Presentation
from pydantic import BaseModel, Field

# -------------------- Environment & Config --------------------
load_dotenv()

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY", "")
MODEL = os.getenv("LLM_MODEL", "gpt-4.1-mini")

KB_DIR = "kb"
DATA_DIR = "data"
INDEX_DIR = f"{DATA_DIR}/index"
LEARNING_DIR = f"{DATA_DIR}/learning"
LOGS_DIR = f"{DATA_DIR}/logs"
CONV_DIR = f"{DATA_DIR}/conversations"
SESSIONS_FILE = f"{DATA_DIR}/sessions.json"

# Optional: additional KB roots (e.g., synced OneDrive/SharePoint folders)
KB_EXTRA_DIRS = [p.strip() for p in os.getenv("KB_EXTRA_DIRS", "").split(",") if p.strip()]

TOP_K = 4  # number of chunks to retrieve

# Ensure directories exist
for d in [DATA_DIR, KB_DIR, INDEX_DIR, LEARNING_DIR, LOGS_DIR, CONV_DIR]:
    Path(d).mkdir(parents=True, exist_ok=True)
for d in KB_EXTRA_DIRS:
    Path(d).mkdir(parents=True, exist_ok=True)

# -------------------- Logging --------------------
logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s")
log = logging.getLogger("comemr-backend")

# -------------------- System Instruction (human-friendly fallback) --------------------
SYSTEM_INSTRUCTION = """
You are ComEMR, a helpful program and clinical support assistant.

CRITICAL BEHAVIOR:
1) Always use the internal Knowledge Base (KB) first.
2) If KB is incomplete, respond naturally using concise, practical guidance; do not say “KB has no answer” or label content as “External information”.
3) Cite KB sources (file names) when you use KB content.
4) Ask a follow-up question ONLY if essential information is missing.
5) Never invent ComEMR policies; if unsure, give general best-practice guidance.
""".strip()

# -------------------- FastAPI App --------------------
app = FastAPI(title="ComEMR Support")

app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "http://localhost:5173", "http://127.0.0.1:5173", "*"
    ],  # tighten in production
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# -------------------- Middleware: Request IDs --------------------
@app.middleware("http")
async def add_request_id(request: Request, call_next):
    rid = str(uuid.uuid4())[:8]
    request.state.request_id = rid
    try:
        response = await call_next(request)
    except Exception as e:
        log.error(f"[{rid}] Unhandled error: {e}")
        return JSONResponse(status_code=500, content={"error": "Unhandled", "detail": str(e), "request_id": rid})
    response.headers["X-Request-ID"] = rid
    return response

def _rid(request: Request) -> str:
    return getattr(request.state, "request_id", "unknown")

# -------------------- Atomic I/O helpers (Windows-safe) --------------------
def ensure_dir(path: Path):
    path.mkdir(parents=True, exist_ok=True)

def atomic_write_json(path: str | Path, obj: Any):
    """Windows-safe atomic JSON write."""
    p = Path(path)
    ensure_dir(p.parent)
    tmp_path = None
    try:
        fd, tmp_path = tempfile.mkstemp(prefix=p.name + ".", suffix=".tmp", dir=p.parent)
        with os.fdopen(fd, "w", encoding="utf-8") as f:
            json.dump(obj, f, indent=2, ensure_ascii=False)
        os.replace(tmp_path, p)  # atomic replace on same filesystem
    except PermissionError:
        # Fallback if Windows locks the tmp
        with open(p, "w", encoding="utf-8") as f:
            json.dump(obj, f, indent=2, ensure_ascii=False)
    finally:
        if tmp_path and os.path.exists(tmp_path):
            try:
                os.remove(tmp_path)
            except Exception:
                pass

def atomic_append_jsonl(path: str | Path, obj: Dict[str, Any]):
    p = Path(path)
    ensure_dir(p.parent)
    with open(p, "a", encoding="utf-8") as f:
        f.write(json.dumps(obj, ensure_ascii=False) + "\n")

# -------------------- Knowledge Base Index --------------------
class KBIndex:
    def __init__(self):
        self.embedder = SentenceTransformer("all-MiniLM-L6-v2")
        self.index: Optional[faiss.IndexFlatIP] = None
        self.texts: List[str] = []
        self.meta: List[Dict[str, Any]] = []

    # --- FAQ helpers -------------------------------------------------
    def _is_faq_file(self, f: str, root: str) -> bool:
        name = str(f).lower()
        path = str(root).lower()
        return ("faq" in name) or ("faq" in path) or ("frequently asked" in name)

    def _split_faq(self, text: str) -> List[str]:
        """
        Extract Q/A blocks from FAQ-like documents.
        Matches lines starting with Q... followed by A... up to next Q or EOF.
        """
        blocks = []
        pattern = r"(?:^|\n)\s*(Q(?:uestion)?[\.:\)]?\s*)(.*?)(?:\r?\n+)\s*(A(?:nswer)?[\.:\)]?\s*)(.*?)(?=\n\s*Q|$)"
        for m in re.finditer(pattern, text, flags=re.IGNORECASE | re.DOTALL):
            q = m.group(2).strip()
            a = m.group(4).strip()
            if q and a:
                blocks.append(f"Q: {q}\nA: {a}")
        return blocks

    # --- loaders for different formats -------------------------------
    def _load_docx(self, path: str) -> str:
        try:
            doc = Document(path)
            return "\n".join(p.text for p in doc.paragraphs)
        except Exception:
            return ""

    def _load_pptx(self, path: str) -> str:
        try:
            prs = Presentation(path)
            lines = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        lines.append(shape.text)
            return "\n".join(lines)
        except Exception:
            return ""

    def _load_odt(self, path: str) -> str:
        # Best-effort extraction from content.xml
        try:
            with zipfile.ZipFile(path) as z:
                xml = z.read("content.xml").decode("utf-8", errors="ignore")
                text = re.sub(r"<[^>]+>", " ", xml)  # strip XML tags
                return re.sub(r"\s+", " ", text).strip()
        except Exception:
            return ""

    def _load_pdf(self, path: str) -> str:
        try:
            reader = PdfReader(path)
            return "\n".join((p.extract_text() or "") for p in reader.pages)
        except Exception:
            return ""

    def _load_text(self, path: str) -> str:
        try:
            return open(path, encoding="utf-8", errors="ignore").read()
        except Exception:
            return ""

    def _load(self, path: str) -> str:
        ext = os.path.splitext(path)[1].lower()
        if ext == ".pdf":
            return self._load_pdf(path)
        elif ext == ".docx":
            return self._load_docx(path)
        elif ext == ".pptx":
            return self._load_pptx(path)
        elif ext == ".odt":
            return self._load_odt(path)
        else:  # Plain text for .txt, .md, .csv, .json, etc.
            return self._load_text(path)

    # --- chunking ----------------------------------------------------
    def _chunk(self, text: str, size: int = 900, overlap: int = 150) -> List[str]:
        out, i = [], 0
        n = len(text)
        if n == 0:
            return out
        step = max(1, size - overlap)
        while i < n:
            out.append(text[i: i + size])
            i += step
        return out

    # --- build / load / search --------------------------------------
    def build(self):
        texts, meta = [], []
        kb_roots = [KB_DIR] + KB_EXTRA_DIRS  # index kb/ + any extra folders

        for base in kb_roots:
            for root, _, files in os.walk(base):
                for f in files:
                    path = os.path.join(root, f)
                    content = self._load(path)
                    if not content:
                        continue  # skip unreadable/empty files

                    faq_chunks = self._is_faq_file(f, root) and self._split_faq(content) or []
                    chunks = faq_chunks if faq_chunks else self._chunk(content)

                    for chunk in chunks:
                        texts.append(chunk)
                        meta.append({
                            "file": f,
                            "path": root,
                            "ext": os.path.splitext(f)[1].lower(),
                            "is_faq": self._is_faq_file(f, root),
                        })

        if not texts:
            self.index = None
            self.texts, self.meta = [], []
            log.warning("[KB] No texts were found during build.")
            return

        emb = self.embedder.encode(texts, normalize_embeddings=True)
        self.index = faiss.IndexFlatIP(emb.shape[1])
        self.index.add(np.array(emb).astype("float32"))
        self.texts, self.meta = texts, meta

        # Persist index + metadata
        try:
            faiss.write_index(self.index, f"{INDEX_DIR}/kb.faiss")
            with open(f"{INDEX_DIR}/kb.pkl", "wb") as pf:
                pickle.dump({"texts": texts, "meta": meta}, pf)
            log.info(f"[KB] Build complete: chunks={len(texts)}")
        except Exception as e:
            log.error(f"[KB] Persist failed: {e}")

    def load(self):
        try:
            self.index = faiss.read_index(f"{INDEX_DIR}/kb.faiss")
            with open(f"{INDEX_DIR}/kb.pkl", "rb") as pf:
                data = pickle.load(pf)
            self.texts, self.meta = data.get("texts", []), data.get("meta", [])
            log.info(f"[KB] Loaded index: chunks={len(self.texts)}")
        except Exception as e:
            log.warning(f"[KB] Load failed ({e}); rebuilding...")
            self.build()

    def search(self, query: str, k: int) -> List[Dict[str, Any]]:
        if not self.index or not self.texts:
            return []
        q = self.embedder.encode([query], normalize_embeddings=True).astype("float32")
        scores, ids = self.index.search(q, k)
        results = []
        for i, s in zip(ids[0], scores[0]):
            if i >= 0:
                results.append({"score": float(s), "text": self.texts[i], "meta": self.meta[i]})
        return results

kb = KBIndex()
kb.load()

# -------------------- Namespace Filter (optional) --------------------
def namespace_allowed(ns: Optional[str], meta: dict) -> bool:
    """
    If ns provided, require it to appear in path or file; else allow.
    """
    ns = (ns or "").strip().lower()
    if not ns:
        return True
    file_name = (meta.get("file") or "").lower()
    folder_path = (meta.get("path") or "").lower()
    return (ns in file_name) or (ns in folder_path)

# -------------------- Prompt Builder (role-agnostic) --------------------
def build_prompt(question: str, kb_hits: List[Dict[str, Any]], ns: Optional[str]) -> str:
    """
    Role-agnostic prompt builder: focuses on KB-first behavior and clarity.
    """
    ctx = "\n\n".join(f"SOURCE ({h['meta']['file']}):\n{h['text']}" for h in kb_hits)
    ns_line = f"NAMESPACE: {ns}" if ns else "NAMESPACE: default"

    return f"""
{SYSTEM_INSTRUCTION}

{ns_line}
COMMUNICATION STYLE: Provide clear, concise, and accurate guidance. Prefer step-by-step instructions when procedural.

INTERNAL KNOWLEDGE BASE:
{ctx if ctx else "No relevant KB content."}

USER QUESTION:
{question}

RESPONSE RULES:
- Prefer KB content when available and cite source filenames
- If KB is insufficient, respond naturally with practical steps (no 'External information' label)
- Ask follow-up ONLY if essential info is missing
""".strip()

# -------------------- Confidence Phrasing --------------------
def confidence_prefix(score: float) -> str:
    if score >= 0.75:
        return ""
    if score >= 0.5:
        return "Based on available information, "
    return "Here’s the best general guidance: "

# -------------------- FAQ Auto-Promotion --------------------
def promote_to_faq(question: str, answer: str, confidence: float):
    if confidence < 0.7:
        return
    faq_path = os.path.join(KB_DIR, "auto_faq.md")
    try:
        with open(faq_path, "a", encoding="utf-8") as f:
            f.write(f"\n## Q: {question}\n{answer}\n")
        kb.build()
    except Exception as e:
        log.warning(f"[FAQ] Promotion failed: {e}")

# -------------------- User Intent Memory --------------------
INTENT_FILE = f"{LEARNING_DIR}/intents.json"

def log_intent(question: str, embedding: np.ndarray):
    intents = []
    if Path(INTENT_FILE).exists():
        try:
            with open(INTENT_FILE, "r", encoding="utf-8") as f:
                intents = json.load(f)
            if not isinstance(intents, list):
                intents = []
        except Exception:
            intents = []
    entry = {"question": question, "embedding": embedding.tolist(), "ts": datetime.utcnow().isoformat()}
    atomic_write_json(INTENT_FILE, intents + [entry])

# -------------------- Answer Quality Logs --------------------
def log_quality(data: Dict[str, Any]):
    atomic_append_jsonl(f"{LOGS_DIR}/answers.jsonl", data)

# -------------------- Conversation Storage --------------------
def _messages_path(session_id: int) -> str:
    return os.path.join(CONV_DIR, f"{session_id}.jsonl")

def load_messages(session_id: int) -> List[Dict[str, Any]]:
    path = _messages_path(session_id)
    msgs = []
    if os.path.exists(path):
        with open(path, "r", encoding="utf-8") as f:
            for line in f:
                try:
                    msgs.append(json.loads(line))
                except Exception:
                    continue
    return msgs

def append_message(session_id: int, role: str, content: str, meta: dict = None) -> Dict[str, Any]:
    rec = {
        "role": role,  # "user" | "assistant" | "system"
        "content": content,
        "meta": meta or {},
        "timestamp": datetime.utcnow().isoformat(),
    }
    Path(CONV_DIR).mkdir(parents=True, exist_ok=True)
    atomic_append_jsonl(_messages_path(session_id), rec)
    return rec

# -------------------- Sessions (safe ID generation) --------------------
def safe_session_id(sid) -> int:
    try:
        return int(sid) if sid is not None and str(sid).strip() != "" else 1
    except Exception:
        return 1

def _load_sessions() -> List[Dict[str, Any]]:
    if os.path.exists(SESSIONS_FILE):
        try:
            with open(SESSIONS_FILE, "r", encoding="utf-8") as f:
                data = json.load(f)
            return data if isinstance(data, list) else []
        except Exception:
            return []
    return []

def _save_sessions(sessions: List[Dict[str, Any]]):
    if not isinstance(sessions, list):
        sessions = []
    Path(DATA_DIR).mkdir(parents=True, exist_ok=True)
    atomic_write_json(SESSIONS_FILE, sessions)

def _next_session_id(sessions: List[Dict[str, Any]]) -> int:
    max_id = 0
    for s in sessions:
        try:
            max_id = max(max_id, int(s.get("id", 0)))
        except Exception:
            continue
    return max_id + 1 if max_id >= 1 else 1

# -------------------- Base Sessions Endpoints (list/create) --------------------
class SessionCreate(BaseModel):
    name: str = Field(..., min_length=1, max_length=200)
    namespace: Optional[str] = Field(default="default", max_length=50)

api = APIRouter()

@api.get("/sessions")
def list_sessions():
    sessions = _load_sessions()
    if not isinstance(sessions, list):
        sessions = []
        _save_sessions(sessions)
    return sessions

@api.post("/sessions")
def create_session(payload: SessionCreate):
    raw_name = (payload.name or "").strip()
    if not raw_name:
        raise HTTPException(status_code=400, detail={"error": "EMPTY_NAME", "message": "Session name cannot be empty."})

    sessions = _load_sessions()
    if not isinstance(sessions, list):
        sessions = []

    new_id = _next_session_id(sessions)
    created_ts = datetime.utcnow().isoformat()

    new_session = {
        "id": new_id,
        "name": raw_name,
        "created": created_ts,
        "namespace": (payload.namespace or "default"),
    }
    sessions.append(new_session)
    _save_sessions(sessions)

    # Ensure the conversation file exists for this session
    Path(CONV_DIR).mkdir(parents=True, exist_ok=True)
    Path(_messages_path(new_id)).touch()

    return new_session

# Mount router at both /sessions and /api/sessions
app.include_router(api)                  # /sessions
app.include_router(api, prefix="/api")   # /api/sessions

# -------------------- Existing Session Operations by ID --------------------
@app.get("/sessions/{sid}")
def get_session(sid: int):
    sessions = _load_sessions()
    s = next((x for x in sessions if x.get("id") == sid), None)
    if not s:
        raise HTTPException(status_code=404, detail="Session not found")
    return s

@app.delete("/sessions/{sid}")
def delete_session(sid: int):
    sessions = _load_sessions()
    idx = next((i for i, s in enumerate(sessions) if s.get("id") == sid), None)
    if idx is None:
        raise HTTPException(status_code=404, detail="Session not found")
    removed = sessions.pop(idx)
    _save_sessions(sessions)
    try:
        os.remove(_messages_path(sid))
    except Exception:
        pass
    return {"deleted": removed}

@app.get("/sessions/{sid}/messages")
def list_session_messages(sid: int):
    sessions = _load_sessions()
    if not any(s.get("id") == sid for s in sessions):
        raise HTTPException(status_code=404, detail="Session not found")
    return load_messages(sid)

@app.post("/sessions/{sid}/messages")
async def add_session_message(
    sid: int,
    request: Request,
    role: Optional[str] = Form(None),
    content: Optional[str] = Form(None),
    payload: Dict[str, Any] = Body(None),
):
    # Accept JSON or form
    if payload:
        role = role or (payload or {}).get("role")
        content = content or (payload or {}).get("content")

    if role not in ("user", "assistant", "system") or not content or not str(content).strip():
        raise HTTPException(status_code=422, detail="role must be user|assistant|system and content required")

    sessions = _load_sessions()
    if not any(s.get("id") == sid for s in sessions):
        raise HTTPException(status_code=404, detail="Session not found")

    rec = append_message(sid, role, content.strip(), meta={})
    return rec

# -------------------- OpenAI Streaming + Non-stream fallback --------------------
def stream_openai(prompt: str) -> Iterator[str]:
    client = OpenAI(api_key=OPENAI_API_KEY)
    try:
        with client.responses.stream(
            model=MODEL,
            input=prompt,
            temperature=0.2,
        ) as stream:
            for event in stream:
                if event.type == "response.output_text.delta":
                    yield event.delta
    except Exception as e:
        log.error(f"Streaming error: {e}")
        # Fallback: non-stream response (still no prompt leaks)
        try:
            resp = client.responses.create(model=MODEL, input=prompt, temperature=0.2)
            full_text = getattr(resp, "output_text", "") or ""
        except Exception as e2:
            log.error(f"Fallback completion error: {e2}")
            full_text = ""
        # Yield chunked to satisfy StreamingResponse contract
        for chunk in re.findall(r".{1,2000}", full_text, flags=re.DOTALL):
            yield chunk

def complete_openai(prompt: str) -> str:
    client = OpenAI(api_key=OPENAI_API_KEY)
    try:
        resp = client.responses.create(model=MODEL, input=prompt, temperature=0.2)
        return getattr(resp, "output_text", "") or ""
    except Exception as e:
        log.error(f"Non-stream completion error: {e}")
        return ""

# -------------------- Common Fallbacks --------------------
COMMON_FALLBACKS = {
    "reset password": "To reset your password, use the Forgot Password option on the login screen. If that fails, contact your supervisor.",
    "comemr devices": "ComEMR runs on program-issued tablets or smartphones. Contact your supervisor if your device is missing or faulty.",
}

# -------------------- Chat Endpoint --------------------
@app.post("/chat")
async def chat(
    request: Request,
    prompt: Optional[str] = Form(None),
    role: Optional[str] = Form(None),               # accepted but unused (role-agnostic)
    session_id: Optional[int] = Form(None),
    ns: Optional[str] = Form(None),
    stream: Optional[bool] = Form(True),            # toggle to get clean JSON response
    # Also support JSON payloads
    json_body: Dict[str, Any] = Body(None),
):
    rid = _rid(request)

    # Merge Form + JSON
    if json_body:
        prompt = prompt or json_body.get("prompt")
        role = (json_body.get("role") or role)
        session_id = json_body.get("session_id", session_id)
        ns = json_body.get("ns", ns)
        stream = json_body.get("stream", stream)

    if not prompt or not str(prompt).strip():
        raise HTTPException(status_code=400, detail={"error": "BadRequest", "detail": "prompt required", "request_id": rid})

    ns = (ns or "").strip().lower()
    session_id = safe_session_id(session_id)

    # Validate or auto-create the session file
    sessions = _load_sessions()
    if not any(s.get("id") == session_id for s in sessions):
        new_session = {
            "id": int(session_id),
            "name": f"Session {session_id}",
            "created": datetime.utcnow().isoformat(),
            "namespace": ns or "default",
        }
        sessions.append(new_session)
        _save_sessions(sessions)
        Path(_messages_path(int(session_id))).touch()

    # -------------------- KB search (role-agnostic) --------------------
    hits = kb.search(prompt, TOP_K)

    # Dynamic threshold: lower for short queries, higher for longer ones
    dyn_thr = 0.35 if len(prompt) < 20 else 0.65
    strong_hits = [h for h in hits if h.get("score", 0.0) >= dyn_thr]

    # Fallback: if no strong hits, include top 2 hits
    if not strong_hits and hits:
        strong_hits = hits[:2]

    # Confidence from strong hits
    confidence = max((h.get("score", 0.0) for h in strong_hits), default=0.0)

    # Score log (debugging)
    log.info(f"[{rid}] KB scores: {[round(h['score'],3) for h in hits]}")

    # Intent memory (non-fatal)
    try:
        emb = kb.embedder.encode(prompt)
        log_intent(prompt, emb)
    except Exception as e:
        log.warning(f"[{rid}] intent logging failed: {e}")

    # Persist the user message (so sidebar shows immediately)
    append_message(int(session_id), "user", prompt, meta={
        "role": (role or "unspecified"),
        "ns": ns or "default",
        "kb_hits": [{"file": h["meta"]["file"], "path": h["meta"].get("path"), "score": h.get("score", 0.0)} for h in strong_hits],
        "request_id": rid,
    })

    composed_prompt = build_prompt(prompt, strong_hits, ns)
    total_hits = len(hits)
    log.info(
        f"[{rid}] chat: ns={ns or 'default'} session={session_id} "
        f"msg_len={len(prompt)} total_hits={total_hits} strong_hits={len(strong_hits)} conf={confidence:.3f}"
    )

    # Common fallbacks: short-circuit for frequent queries
    for k, v in COMMON_FALLBACKS.items():
        if k in prompt.lower():
            answer = v
            # Persist assistant message
            append_message(int(session_id), "assistant", answer, meta={
                "role": (role or "unspecified"),
                "ns": ns or "default",
                "confidence": confidence,
                "sources": [h["meta"]["file"] for h in strong_hits],
                "timestamp": datetime.utcnow().isoformat(),
                "request_id": rid,
            })
            return {
                "response": answer,
                "session_id": session_id,
                "ns": ns or "default",
                "role": (role or "unspecified"),
                "confidence": confidence,
                "sources": [h["meta"]["file"] for h in strong_hits],
                "request_id": rid,
            }

    # --- Non-stream (clean JSON response) ---
    if not stream:
        answer = complete_openai(composed_prompt).strip()
        prefix = confidence_prefix(confidence)
        final_answer = (prefix + answer) if answer else "Samahani, naomba ufafanuzi zaidi ili nikusaidie vizuri."

        # Persist assistant message
        append_message(int(session_id), "assistant", final_answer, meta={
            "role": (role or "unspecified"),
            "ns": ns or "default",
            "confidence": confidence,
            "sources": [h["meta"]["file"] for h in strong_hits],
            "timestamp": datetime.utcnow().isoformat(),
            "request_id": rid,
        })

        # Quality log
        log_quality({
            "question": prompt,
            "role": (role or "unspecified"),
            "ns": ns or "default",
            "confidence": confidence,
            "timestamp": datetime.utcnow().isoformat(),
            "request_id": rid,
        })

        # FAQ auto-promotion if high confidence
        promote_to_faq(prompt, final_answer, confidence)

        # Follow-up prompt if weak
        if confidence < 0.4:
            final_answer += "\n\nCould you clarify or provide more details so I can help better?"

        # Clean JSON response (no prompt leaks)
        return {
            "response": final_answer.strip(),
            "session_id": session_id,
            "ns": ns or "default",
            "role": (role or "unspecified"),
            "confidence": confidence,
            "sources": [h["meta"]["file"] for h in strong_hits],
            "request_id": rid,
        }

    # --- Streaming (text/plain) ---
    def stream_gen():
        prefix = confidence_prefix(confidence)
        # Prepend prefix once at start of stream, then stream tokens
        yield prefix
        output = ""
        try:
            for token in stream_openai(composed_prompt):
                output += token
                yield token
        except Exception as e:
            log.error(f"[{rid}] Streaming failed: {e}")
            yield "\n[Error] Assistant temporarily unavailable."

        # Persist assistant final text with sources/confidence after stream completes
        final_text = prefix + output
        append_message(int(session_id), "assistant", final_text, meta={
            "role": (role or "unspecified"),
            "ns": ns or "default",
            "confidence": confidence,
            "sources": [h["meta"]["file"] for h in strong_hits],
            "timestamp": datetime.utcnow().isoformat(),
            "request_id": rid,
        })

        # Quality log
        log_quality({
            "question": prompt,
            "role": (role or "unspecified"),
            "ns": ns or "default",
            "confidence": confidence,
            "timestamp": datetime.utcnow().isoformat(),
            "request_id": rid,
        })

        # FAQ auto-promotion if high confidence
        promote_to_faq(prompt, final_text, confidence)

    return StreamingResponse(stream_gen(), media_type="text/plain")

# -------------------- ChatGPT-style UI --------------------
@app.get("/", response_class=HTMLResponse)
def ui():
    path = "chat_ui.html"
    if os.path.exists(path):
        return open(path, "r", encoding="utf-8").read()
    return "<h3>ComEMR Support Assistant</h3>"

# -------------------- Health + Config --------------------
@app.get("/health")
def health():
    return {
        "status": "ok",
        "kb_chunks": len(kb.texts),
        "model": MODEL,
    }

@app.get("/config")
def config():
    index_exists = os.path.exists(f"{INDEX_DIR}/kb.faiss") and os.path.exists(f"{INDEX_DIR}/kb.pkl")
    return {
        "provider": "openai",
        "model": MODEL,
        "kb_default_dir": KB_DIR,
        "kb_extra_dirs": KB_EXTRA_DIRS,
        "kb_index_exists": index_exists,
        "kb_chunks": len(kb.texts),
        "sessions_file": SESSIONS_FILE,
        "conversations_dir": CONV_DIR,
    }

# -------------------- Admin Learning Approval --------------------
@app.post("/admin/approve")
async def approve_learning(
    request: Request,
    question: Optional[str] = Form(None),
    answer: Optional[str] = Form(None),
    role: Optional[str] = Form(None),   # accepted for metadata only
    payload: Dict[str, Any] = Body(None),
):
    rid = _rid(request)

    # Accept both JSON and form
    if payload:
        question = question or payload.get("question")
        answer = answer or payload.get("answer")
        role = role or payload.get("role")

    if not (question and answer):
        raise HTTPException(status_code=400, detail={"error": "BadRequest", "detail": "question and answer required", "request_id": rid})

    fname = "approved_general.md"  # role-agnostic approval file
    try:
        with open(os.path.join(KB_DIR, fname), "a", encoding="utf-8") as f:
            f.write(f"\n## Q: {question}\n{answer}\n")
    except Exception as e:
        log.error(f"[{rid}] approve write failed: {e}")
        raise HTTPException(status_code=500, detail={"error": "WriteFailed", "detail": "Could not persist approval", "request_id": rid})

    kb.build()
    return {"status": "added_to_kb", "file": fname}

# -------------------- Admin KB Reindex --------------------
@app.post("/admin/reindex")
def reindex():
    """
    Rebuild the entire Knowledge Base index.
    """
    kb.build()
    return {
        "status": "rebuilt",
        "kb_chunks": len(getattr(kb, "texts", []))
    }

# -------------------- KB Stats --------------------
@app.get("/kb/stats")
def kb_stats():
    """
    Return Knowledge Base statistics.
    """
    try:
        total_chunks = len(getattr(kb, "texts", []))
        return {
            "namespace": "default",
            "chunks": total_chunks
        }
    except Exception as e:
        log.error(f"KB stats error: {e}")
        raise HTTPException(
            status_code=500,
            detail="Could not retrieve KB stats"
        )
